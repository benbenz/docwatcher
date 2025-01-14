from urllib.parse import urlparse
from crawler.handlers import get_filename , get_content_type , LocalStorageHandler 
from crawler.core import CrawlerMode, bcolors , FileStatus 
from crawler.helper import clean_url
# .pdf
#from PyPDF4 import PdfFileReader
from pypdf import PdfReader
# .doc , .docx
from docx import Document as WordDocument
# .ppt , .pptx , .pptm
from pptx import Presentation
# .rtf
from striprtf.striprtf import rtf_to_text
# .html
from bs4 import BeautifulSoup
# PDF image treatment
from PIL import Image

import os
import re
import csv
import uuid
import traceback
import subprocess
import json
from email.utils import parsedate_to_datetime
from datetime import datetime, timedelta
from django.utils.timezone import make_aware

import logging
logger = logging.getLogger("DocCrawler")

# load django stuff
# MAKE SURE ROOT/www is also in the PYTHONPATH !!!
os.environ['DJANGO_SETTINGS_MODULE'] = 'docwatcher.settings'
from django.core.wsgi import get_wsgi_application
application = get_wsgi_application()
# now we can load the model :)
# note that we don't have www.docs.models (because of PYTHONPATH)
from docs.models import Document
from django.db.models import Q
from django.db import models

def get_doctype_by_content(response):
    try:
        content = response.text.strip()
        if "<!doctype html>" in content[:64]:
            return Document.DocumentType.HTML
        elif "%PDF-" in content[:32]:
            return Document.DocumentType.PDF
        elif "word/" in content[:32]:
            return Document.DocumentType.DOCX
    except:
        pass
    return Document.DocumentType.UNKNOWN

def get_doctype_by_extension(file_extension):
    if file_extension==".pdf" :
        return Document.DocumentType.PDF
    elif file_extension==".html":
        return Document.DocumentType.HTML
    elif file_extension=="text/plain":
        return Document.DocumentType.TEXT
    elif file_extension=="application/rtf":
        return Document.DocumentType.RTF
    elif file_extension=="application/msword":
        return Document.DocumentType.DOC
    elif file_extension==".docx":
        return Document.DocumentType.DOCX
    elif file_extension==".ppt":
        return Document.DocumentType.PPT
    elif file_extension==".pptx":
        return Document.DocumentType.PPTX
    elif file_extension==".pptm":
        return Document.DocumentType.PPTM   
    return Document.DocumentType.UNKNOWN    

def get_doctype_by_url(response):
    try:
        last_part = response.url.rsplit('/', 1)[-1]
        file_name, file_extension = os.path.splitext(last_part)
        if file_extension:
            file_extension = file_extension.lower()
            return get_doctype_by_extension(file_extension)
    except:
        pass
    return Document.DocumentType.UNKNOWN        

def get_doctype(response,try_all_methods=False):
    content_type = get_content_type(response)
    if content_type=="application/pdf" :
        return Document.DocumentType.PDF
    elif content_type=="text/html":
        return Document.DocumentType.HTML
    elif content_type=="text/plain":
        return Document.DocumentType.TEXT
    elif content_type=="application/rtf":
        return Document.DocumentType.RTF
    elif content_type=="application/msword":
        return Document.DocumentType.DOC
    elif content_type=="application/vnd.openxmlformats" or content_type=="officedocument.wordprocessingml.document":
        return Document.DocumentType.DOCX
    elif content_type=="application/vnd.ms-powerpoint":
        return Document.DocumentType.PPT
    elif content_type=="application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return Document.DocumentType.PPTX
    elif content_type=="application/vnd.ms-powerpoint.presentation.macroEnabled.12":
        return Document.DocumentType.PPTM

    if try_all_methods:
        doctype = get_doctype_by_content(response)
        if doctype != Document.DocumentType.UNKNOWN:
            return doctype
        doctype = get_doctype_by_url(response)
        if doctype != Document.DocumentType.UNKNOWN:
            return doctype

    return Document.DocumentType.UNKNOWN

def reset_eof_of_pdf_return_stream(pdf_stream_in:list):
    # find the line position of the EOF
    found_it = False
    for i, x in enumerate(pdf_stream_in[::-1]):
        if b'%%EOF' in x:
            actual_line = len(pdf_stream_in)-i
            logger.debug(f'EOF found at line position {-i} = actual {actual_line}, with value {x}')
            found_it = True
            break

    # return the list up to that point
    if found_it:
        return pdf_stream_in[:actual_line]   
    else:
        pdf_stream_in.append(b'%%EOF')
        return pdf_stream_in

def recover_PDF(path):
    with open(path, 'rb') as p:
        txt = (p.readlines())
        # get the new list terminating correctly
        txtx = reset_eof_of_pdf_return_stream(txt)
        # write to new pdf
        tmp_file = str(uuid.uuid4())[:8] + ".pdf"
        with open(tmp_file, 'wb') as tmpf:
            tmpf.writelines(txtx)
        try:
            pdf = PdfReader(tmp_file) #PdfFileReader(tmp_file)
            os.remove(tmp_file)
            return pdf
        except Exception as e:
            os.remove(tmp_file)
            raise e
    return None

def get_header_http_last_modified(response):
    last_modified = response.headers.get('Last-Modified') or response.headers.get('last-modified')
    if last_modified:
        last_modified = parsedate_to_datetime(last_modified)
    return last_modified

def get_header_http_length(response):
    return response.headers.get('Content-Length') or response.headers.get('content-length') or -1    

def get_header_http_encoding(response):
    return response.headers.get('Content-Encoding') or response.headers.get('content-encoding') or '' 

class AllInOneHandler(LocalStorageHandler):

    def __init__(self, directory, subdirectory, use_ocr=None):
        #super().__init__(directory,subdirectory)
        super().__init__(directory,None) # we will dynamically use the netloc for the subdirectory
        self.using_ocr = False
        if use_ocr==False:
            logger.warning("Specifically NOT using OCR")
            self.process_PDF_body = self.process_PDF_body_NO_OCR
        else:
            try:
                import easyocr
                logger.info_plus("using OCR")
                self.process_PDF_body = self.process_PDF_body_with_OCR
                self.using_ocr = True
            except (ImportError,ModuleNotFoundError) as e:
                logger.warning("NOT using OCR ({0})".format(e))
                self.process_PDF_body = self.process_PDF_body_NO_OCR
                #traceback.print_exc()

    def process_PDF_body_NO_OCR(self,url,path,pdf,sleep=0):
        needs_ocr = False
        
        body = "\n".join([p.extract_text() for p in pdf.pages])
        
        if body == '':
            logger.warning("file needs OCR")
            needs_ocr = True        

        return body , needs_ocr , False

    def process_PDF_page_with_OCR(self,url,path,page,page_count,ocr_reader,sleep=0):
        
        debug = True
        if debug:
            logger.debug("processing page {0}".format(page_count))
        page_body = page.extract_text()
        img_count = 0
        rotation  = page.get('/Rotate')
        found_extra_text = False
        file_root = str(uuid.uuid4())[:8]

        try:
            for image in page.images:
                if self.do_stop:
                    return page_body , found_extra_text
                if debug:
                    logger.debug("processing page {0} image {1} for {2}".format(page_count,img_count,url))
                with open(image.name, "wb") as fp:
                    fp.write(image.data)
                #im0 = Image.open(image.name)
                best_text = None
                # spawn the process out (it can crash)
                try:
                    if sleep!=0:
                        process_args = ['python','docspider/ocr.py','-s',sleep,image.name]
                    else:
                        process_args = ['python','docspider/ocr.py',image.name]
                    process = subprocess.run(process_args,capture_output=True)
                    stdout  = process.stdout
                    stderr  = process.stderr
                    ex_code = process.returncode
                    in_data = False
                    for line in stdout.split(b'\n'):
                        if line.startswith(b'RESULT='):
                            line=line.replace(b'RESULT=',b'')
                            byte_array = bytearray.fromhex(line.decode())
                            str_dump = byte_array.decode('utf-8')
                            json_result = json.loads(str_dump)
                            best_text = json_result.get('best_text')
                            break
                    if ex_code != 0:
                        logger.warning("Error processing image #{0} page #{1} for {2}: exit code = {3}\nstdout={4}\nstderr={5}".format(img_count,page_count,url,ex_code,stdout,stderr))
                except Exception as e:
                    logger.error("Error running process {0} {1} {2} {3}".format(process_args,stdout,stderr,best_text))
                    logger.exception(e,exc_info=True)
                    #traceback.print_exc()  
                    try:
                        os.remove(image.name)
                    except:
                        pass
                    continue

                if best_text is not None and best_text:
                    if debug:
                        logger.debug("Found text: {0}".format(best_text))
                    page_body += '\n' + best_text 
                try:
                    os.remove(image.name)
                except:
                    pass
                img_count += 1
        except ValueError as verr:
            msg = str(verr)
            if "not enough image data" in msg:
                logger.warning("Error retrieving image data")
            else:
                raise verr
        
        return page_body , found_extra_text

    def process_PDF_body_with_OCR(self,url,path,pdf,sleep=0):

        default_body = "\n".join([p.extract_text() for p in pdf.pages])
                    
        # 1) https://github.com/JaidedAI/EasyOCR
        # 2) https://github.com/PaddlePaddle/PaddleOCR
        # 3) https://github.com/madmaze/pytesseract

        #import easyocr 
        #ocr_reader = easyocr.Reader(['fr']) 
        ocr_reader = None

        found_extra_text = False

        # https://stackoverflow.com/questions/63983531/use-tesseract-ocr-to-extract-text-from-a-scanned-pdf-folders
        page_count = 0
        body = ''

        logger.info("processing {0} ({1})".format(url,path))

        for page in pdf.pages: 
            try:
                page_body , has_extra_text = self.process_PDF_page_with_OCR(url,path,page,page_count,ocr_reader,sleep)
                if page_body:
                    body += page_body + '\n'
                found_extra_text = has_extra_text or found_extra_text
            except:
                try:
                    page_body = page.extract_text()
                    if page_body:
                        body += page_body + '\n'
                except:
                    pass
            # lets try to clean the images if it crashed
            try:
                for image in page.images:
                    try:
                        os.remove(image.name)
                    except:
                        pass  
                    try:
                        t_img_name = "t"+str(image.name)+".png"
                        os.remove(image.name)
                    except:
                        pass
                if self.do_stop:
                    if not found_extra_text:
                        return default_body , False , False #has_ocr = False because we're not done
                    else:
                        return body , False , False #has_ocr = False because we're not done
            except:
                pass
            page_count += 1
            
        if not found_extra_text:
            return default_body , False , True
        else:
            return body , False , True

    def get_documents(self,doc_types=None,doc_types_exclude=None,for_ocr=False):
        docs = Document.objects.filter(is_handled=True)
        if doc_types:
            docs = docs.filter(doc_type__in=doc_types)
        if doc_types_exclude:
            docs = docs.exclude(doc_type__in=doc_types_exclude)
        if for_ocr == True:
            docs = docs.filter(has_ocr=False)
        return docs

    def update_document(self,the_doc,sleep=0):
        title , body , num_pages , needs_ocr , has_error  , has_ocr = self.process_document(the_doc.url,the_doc.local_file,the_doc.doc_type,None,sleep=sleep)
        if title is not None:
            the_doc.title = title
        if body is not None:
            the_doc.body  = body 
        if num_pages != -1:
            the_doc.num_pages = num_pages
        the_doc.needs_ocr = needs_ocr
        the_doc.has_error = has_error
        the_doc.has_ocr   = has_ocr
        the_doc.save() 

    def process_document(self,url,path,doc_type,response_body=None,sleep=0):
        title         = None
        body          = None
        num_pages     = -1
        needs_ocr     = False
        has_error     = False
        has_ocr       = False

        if doc_type == Document.DocumentType.PDF:
            try:
                with open(path, 'rb') as f:
                    pdf         = PdfReader(f) #PdfFileReader(f)
                    information = pdf.metadata #pdf.getDocumentInfo()
                    num_pages   = len(pdf.pages) #pdf.getNumPages()
                    title       = information.title if information and information.title else None
                    body , needs_ocr , has_ocr = self.process_PDF_body(url,path,pdf,sleep=sleep)
            except Exception as e:
                msg = str(e)
                if "EOF" in msg:
                    try:
                        pdf         = recover_PDF(path)
                        information = pdf.metadata #pdf.getDocumentInfo()
                        num_pages   = len(pdf.pages) #pdf.getNumPages()
                        title       = information.title if information and information.title else None
                        body , needs_ocr , has_ocr = self.process_PDF_body(url,path,pdf,sleep=sleep)
                    except Exception as e2:
                        has_error = True
                        logger.error("ERROR recovering file {0} {1} {2}".format(url,path,e2))
                        #traceback.print_exc()

                else:
                    has_error = True
                    logger.error("ERROR processing file {0} {1} {2}".format(url,path,e))
                    traceback.print_exc()

        elif doc_type in [Document.DocumentType.DOC , Document.DocumentType.DOCX]:
            try:
                with open(path,'rb') as f:
                    worddoc  = WordDocument(f)
                    body     = "\n".join([p.text for p in worddoc.paragraphs])
            except Exception as e:
                has_error = True
                logger.error("ERROR processing file {0} {1} {2}".format(url,path,e))
                traceback.print_exc()

        elif doc_type in [Document.DocumentType.PPT , Document.DocumentType.PPTX , Document.DocumentType.PPTM]:
            try:
                with open(path,'rb') as f:
                    prez  = Presentation(f)
                    body  = ""
                    for slide in prez.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    body += shape.text + '\n'
            except Exception as e:
                has_error = True
                logger.error("ERROR processing file {0} {1} {2}".format(url,path,e))
                traceback.print_exc()
        
        elif doc_type == Document.DocumentType.RTF:
            if response_body:
                body = rtf_to_text(response_body)
            else:
                with open(path,'rb') as f:
                    body = rtf_to_text(f.read())

        elif doc_type == Document.DocumentType.HTML:
            try:
                if response_body:
                    body = response_body
                else:
                    with open(path,'rb') as f:
                        body = f.read()
                soup  = BeautifulSoup(body,'html.parser')
                body  = soup.get_text()
                if soup.title:
                    title = soup.title.string
            except Exception as e:
                has_error = True
                logger.error("ERROR processing file {0} {1} {2}".format(url,path,e))
                traceback.print_exc()  

        return title , body , num_pages , needs_ocr , has_error , has_ocr      

    def process_response(self,path,response):
        parsed_url    = urlparse(response.url)
        domain_name   = parsed_url.netloc
        filename      = get_filename(parsed_url,response)
        doc_type      = get_doctype(response)
        title         = filename
        body          = response.content

        nu_title , nu_body , num_pages , needs_ocr , has_error , has_ocr = self.process_document(response.url,path,doc_type,body)        
        if nu_title is not None:
            title = nu_title
        if nu_body is not None:
            body = nu_body

        last_modified = get_header_http_last_modified(response)

        # cleanup extraneous spaces
        if body:
            if not isinstance(body,bytes) and not isinstance(body,bytearray):
                body = re.sub(r'([\s]{10,}|[\n]{3,})','\n\n',body)

        return domain_name , filename , doc_type , title , body , num_pages, needs_ocr , has_error , has_ocr , last_modified


    def handle(self, response, depth, previous_url, previous_id, *args, **kwargs):
        path , file_status , path_as_id = super().handle(response,*args,**kwargs)

        # the file already existed
        if file_status & FileStatus.EXISTING :
            try:
                # local_file path is unique
                the_doc = Document.objects.get(local_file=path)

                # let's also update the content to match the file (unless it's exactly the same file)
                if file_status & FileStatus.EXACT == 0:
                    domain_name , filename , doc_type , title , body , num_pages , needs_ocr , has_error , has_ocr , last_modified = self.process_response(path,response)
                    the_doc.body  = body 
                    the_doc.title = title
                    the_doc.num_pages = num_pages
                    the_doc.needs_ocr = needs_ocr
                    the_doc.has_error = has_error
                    the_doc.has_ocr   = has_ocr
                    the_doc.last_modified = last_modified
                    the_doc.is_handled = True
                    the_doc.save() 
                    logger.debug("updated existing document entry {0} / {1}".format(the_doc.id,path,the_doc.url))
                return path , file_status , the_doc.id
            except:
                logger.error("INTERNAL Error: an existing file is not registered in the database!")
                return path , file_status , None

        domain_name , filename , doc_type , title , body , num_pages , needs_ocr , has_error , has_ocr , last_modified = self.process_response(path,response)

        final_url = kwargs.get('final_url')
        if final_url:
            final_url = clean_url(final_url) # just to make really sure
        else:
            final_url = clean_url(response.url) 
        if kwargs.get('url'):
            url = clean_url(kwargs.get('url'))
        else:
            url = None

        doc = Document(
            domain      = domain_name , 
            url         = url , 
            final_url   = final_url ,
#            referer     = previous_url or '' ,
#            depth       = depth ,
#            record_date = AUTO
            remote_name = filename ,
            http_length =  get_header_http_length(response),
            http_encoding = get_header_http_encoding(response),
            http_last_modified = get_header_http_last_modified(response),
            http_content_type = get_content_type(response),

            # Content: HTML/PDF + file
            doc_type    = doc_type ,
            title       = title or filename,
            body        = body ,
            size        = len(response.content) ,
            num_pages   = num_pages ,
            needs_ocr   = needs_ocr ,
            has_ocr     = has_ocr ,
            has_error   = has_error ,
            local_file  = path , 
            file_status = file_status ,
            of_interest = False ,
            is_handled  = True ,
            record_date = make_aware(datetime.now())
        )

        # save the new entry
        doc.save()

        # add the referer
        added = None
        if previous_id is not None:
            added = False
            try:
                docreferer = Document.objects.get(pk=previous_id)
                if not doc.referers.contains(docreferer):
                    doc.referers.add(docreferer)
                    added = True
            except:
                pass
        elif previous_url is not None and added==False:
            try:
                docreferer = Document.objects.filter(url=previous_url).latest('record_date')
                if not doc.referers.contains(docreferer):
                    doc.referers.add(docreferer)
                    added = True
            except:
                pass

        # save the relationship
        if added:
           doc.save()

        return path , file_status , doc.id

class DBStatsHandler:

    def __init__(self,domain):
        self.domain = domain

    def get_handled_list(self,crawler_mode):
        list_handled = []
        
        if crawler_mode & CrawlerMode.CRAWL_FULL:
        
            pass

        elif crawler_mode & CrawlerMode.CRAWL_THRU:
            # we're gonna consider that really old non-html documents won't change anymore (2y+)
            # we can add those documents to the 'handled_list' and avoid head() or even get() requests
            # we have to discard HTML documents because we want to crawl them again
            # unless they are very very old as well (3 years+) ! 
            # this is all to minimize the footprint on the server....
            date_today = datetime.today()
            date_html  = make_aware( date_today - timedelta(days=3*365) ) # 3 years
            date_other = make_aware( date_today - timedelta(days=3*365) ) # 3 years 
            q_html  = Q(doc_type=Document.DocumentType.HTML)  & Q(http_last_modified__isnull=False) & Q(http_last_modified__lte=date_html)
            q_other = ~Q(doc_type=Document.DocumentType.HTML) & Q(http_last_modified__isnull=False) & Q(http_last_modified__lte=date_other)
            query   = Q(domain=self.domain,is_handled=True) & (q_html | q_other)
            if self.domain:
                queryset = Document.objects.filter(query)
                #logger.debug(queryset.query)
                for doc in queryset:
                    list_handled.append(doc.url if isinstance(doc.url,str) else doc.url.name)

        elif crawler_mode & CrawlerMode.CRAWL_LIGHT or crawler_mode & CrawlerMode.CRAWL_ULTRA_LIGHT :
            # we're gonna consider more documents as being done ...
            date_today = datetime.today()
            years      = 2 if crawler_mode & CrawlerMode.CRAWL_LIGHT else 1
            date_html  = make_aware( date_today - timedelta(days=years*365) ) 
            date_other = make_aware( date_today - timedelta(days=years*365) ) 
            q_html  = Q(doc_type=Document.DocumentType.HTML)  & Q(http_last_modified__isnull=False) & Q(http_last_modified__lte=date_html)
            q_other = ~Q(doc_type=Document.DocumentType.HTML) & Q(http_last_modified__isnull=False) & Q(http_last_modified__lte=date_other)
            query   = Q(domain=self.domain,is_handled=True) & (q_html | q_other)
            if self.domain:
                queryset = Document.objects.filter(query)
                for doc in queryset:
                    list_handled.append(doc.url if isinstance(doc.url,str) else doc.url.name)

        return list_handled

    def handle(self, response, depth, previous_url, local_name, *args, **kwargs):
        # the job has already been done by the storage handler
        pass

    def get_filenames(self,url,final_url=None):
        result = []
        for doc in Document.objects.filter( Q(is_handled=True) & (Q(url=url) | Q(final_url=final_url)) ):
            result.append(doc.local_file)
        return result

    def find(self,url,response):
        url                = clean_url(url)
        final_url          = clean_url(response.url)
        http_length        = get_header_http_length(response)
        http_encoding      = get_header_http_encoding(response)
        http_last_modified = get_header_http_last_modified(response)
        q_url              = Q(url=url) | Q(final_url=final_url)
        if isinstance(http_length,str):
            try:
                http_length = int(http_length)
            except:
                pass
        try:
            result = None
            if http_last_modified:
                logger.debug("finding with last_modified {0}".format(http_last_modified))
                result = Document.objects.filter(is_handled=True).filter(q_url & Q(http_last_modified=http_last_modified,http_last_modified__isnull=False)).latest('record_date')
            else:
                logger.debug("finding with length={0} and encoding='{1}'".format(http_length,http_encoding))
                # we want the most-recently fetched document to be the same as the one we're comparing it to
                result = Document.objects.filter(is_handled=True).filter(q_url).latest('record_date') # most recently fetched document
                if result.http_length!=http_length or result.http_encoding!=http_encoding: # should be the same in size
                    return None
            return result.id
        except Document.DoesNotExist:
            logger.debug("no document found")
            pass
        except:
            pass
        return None

    def find_recent(self,url,days=7):
        try:
            url = clean_url(url)
            date_today  = datetime.today()
            date_filter = make_aware( date_today - timedelta(days=days) )
            q_http_last_modified = Q(http_last_modified__gte=date_filter) & Q(http_last_modified__isnull=False)
            q_record_date        = Q(record_date__gte=date_filter) & Q(record_date__isnull=False)
            q_url                = Q(url=url) | Q(final_url=url,final_url__isnull=False)
            results = Document.objects.filter(is_handled=True).filter(q_url & (q_record_date | q_http_last_modified))
            #number  = results.count()
            result  = results.latest('record_date')
            return result.id , result.http_content_type , result.http_last_modified , result.record_date
        except Document.DoesNotExist:
            pass
        except Exception as e:
            logger.error("error while finding recent element {0}".format(e))
        return None , None , None , None

    def find_latest(self,url):
        try:
            url   = clean_url(url)
            q_url = Q(url=url) | Q(final_url=url,final_url__isnull=False)
            results = Document.objects.filter(is_handled=True).filter(q_url)
            #number  = results.count()
            result  = results.latest('record_date')
            return result.id , result.http_content_type , result.http_last_modified , result.record_date
        except Document.DoesNotExist:
            pass
        except Exception as e:
            logger.error("error while finding recent element {0}".format(e))
        return None , None , None , None    

    def get_urls_by_referer(self,referer_url,objid=None):
        result = []

        # THIS IS USING a separate cache of URLs... (cf. pre_record_document vs pre_record_document_OLD)
        # queryset = None
        # if objid is not None:
        #     queryset = RecLinkedUrl.objects.filter(referer_id=objid)
        # else:
        #     queryset = RecLinkedUrl.objects.filter(Q(referer__url=referer_clean)|Q(referer__final_url=referer_clean))
        # for reclink in queryset:
        #     result.append({"url":reclink.url,"follow":True})
        # return result

        result = []

        queryset = None
        if objid is not None:
            try:
                queryset = Document.objects.get(pk=objid).links.all()
            except Document.DoesNotExist:
                pass
        if queryset is None:
            referer_clean = clean_url(referer)
            queryset = Document.objects.filter(referers__url=referer_clean)

        for doc in queryset:
            result.append({'url':doc.url,'follow':True})
        return result


    def get_urls_of_interest(self):
        result = set()

        queryset = Document.objects.filter(Q(of_interest=True)|Q(links__of_interest=True)).filter(domain=self.domain).distinct()
        #queryset = Document.objects.filter(Q(of_interest=True)|Q(referer_docs__link__of_interest=True)).filter(domain=self.domain,is_handled=True).distinct()

        for doc in queryset:
            result.add(doc.url)
        return result   

    # def pre_record_clear(self, previous_id, depth):
    #     if not previous_id:
    #         return
    #     RecLinkedUrl.objects.filter(referer_id=previous_id).delete()

    # def pre_record_document(self, previous_id , url):
    #     if not previous_id:
    #         return 
    #     try:
    #         urlobj = RecLinkedUrl.objects.get(url=url,referer_id=previous_id)
    #     except RecLinkedUrl.DoesNotExist:
    #         urlobj = RecLinkedUrl(url=url,referer_id=previous_id)
    #         urlobj.save()

